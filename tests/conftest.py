from __future__ import annotations

import io
from pathlib import Path

import pytest


def make_pdf_bytes(text: str = "Hello, PDF!") -> bytes:
    import fitz

    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((50, 50), text)
    buf = io.BytesIO()
    doc.save(buf)
    doc.close()
    return buf.getvalue()


def make_docx_bytes(text: str = "Hello, DOCX!") -> bytes:
    from docx import Document

    doc = Document()
    doc.add_paragraph(text)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def make_pptx_bytes(text: str = "Hello, PPTX!") -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = top = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, Inches(6), Inches(1))
    tf = txBox.text_frame
    tf.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@pytest.fixture
def pdf_bytes() -> bytes:
    return make_pdf_bytes()


@pytest.fixture
def docx_bytes() -> bytes:
    return make_docx_bytes()


@pytest.fixture
def pptx_bytes() -> bytes:
    return make_pptx_bytes()
